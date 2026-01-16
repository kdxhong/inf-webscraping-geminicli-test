import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tabulate import tabulate
from loguru import logger

def clean_text(text):
    """텍스트 정제: 줄바꿈 제거, 파이프 이스케이프, 공백 제거"""
    if not text:
        return ""
    return text.replace('\n', ' ').replace('|', r'\|').strip()

def convert_pptx_to_md(pptx_path, output_md_path):
    """
    PPTX 파일을 읽어 Markdown 파일로 변환하고 이미지를 추출합니다.
    진짜 표(TABLE) 뿐만 아니라 그리드 형태로 배치된 텍스트 박스도 표로 변환 지원합니다.
    """
    
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    image_dir = os.path.join(os.path.dirname(output_md_path), "images")
    
    if not os.path.exists(image_dir):
        os.makedirs(image_dir)

    logger.info(f"변환 시작: {pptx_path}")
    prs = Presentation(pptx_path)
    
    with open(output_md_path, 'w', encoding='utf-8') as f:
        f.write(f"# {base_name}\n\n")

        for i, slide in enumerate(prs.slides):
            slide_number = i + 1
            f.write(f"---\n\n## Slide {slide_number}\n\n")
            logger.debug(f"Slide {slide_number} 처리 중...")
            
            # 슬라이드 내 모든 모양(Shape) 정보를 수집
            all_shapes = []
            for shape in slide.shapes:
                # 1. 이미지 우선 처리 (이미지는 즉시 출력 또는 위치 저장)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_filename = f"slide{slide_number}_{shape.shape_id}.{image.ext}"
                    image_path = os.path.join(image_dir, image_filename)
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image.blob)
                    rel_path = os.path.join("images", image_filename)
                    all_shapes.append({'type': 'image', 'top': shape.top, 'left': shape.left, 'content': f"![Image]({rel_path})"})
                
                # 2. 진짜 표 처리
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    rows_data = []
                    for row in shape.table.rows:
                        cells_text = [clean_text(cell.text_frame.text) for cell in row.cells]
                        rows_data.append(cells_text)
                    if rows_data:
                        # 수동 표 생성 (표준 호환성 극대화)
                        headers = rows_data[0]
                        body = rows_data[1:]
                        hdr_str = "| " + " | ".join(headers) + " |"
                        sep_str = "| " + " | ".join([":---"] * len(headers)) + " |"
                        table_md = hdr_str + "\n" + sep_str
                        for r_data in body:
                            table_md += "\n| " + " | ".join(r_data) + " |"
                        all_shapes.append({'type': 'table', 'top': shape.top, 'left': shape.left, 'content': table_md})

                # 3. 텍스트 프레임 처리 (가짜 표 감지를 위해 위치 정보 보존)
                elif shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        all_shapes.append({'type': 'text', 'top': shape.top, 'left': shape.left, 'content': text, 'height': shape.height})

            # 수집된 모양들을 위치 기반으로 정렬 및 그룹화 (그리드 감지)
            sorted_shapes = sorted(all_shapes, key=lambda x: (x['top'], x['left']))
            
            # 1단계: 행(Row) 단위로 그룹화
            rows = []
            current_row = []
            last_top = -1e9
            
            for s in sorted_shapes:
                if s['type'] == 'text':
                    if abs(s['top'] - last_top) < 100000: # 같은 행으로 간주
                        current_row.append(s)
                    else:
                        if current_row:
                            rows.append({'type': 'logical_row', 'elements': current_row, 'top': current_row[0]['top']})
                        current_row = [s]
                        last_top = s['top']
                else:
                    if current_row:
                        rows.append({'type': 'logical_row', 'elements': current_row, 'top': current_row[0]['top']})
                        current_row = []
                    rows.append({'type': s['type'], 'content': s['content'], 'top': s['top']})
            if current_row:
                rows.append({'type': 'logical_row', 'elements': current_row, 'top': current_row[0]['top']})

            # 2단계: 연속된 'logical_row'들을 하나의 표로 합치기
            processed_elements = []
            idx = 0
            while idx < len(rows):
                item = rows[idx]
                # 논리적 행이고 요소가 여러 개면 표의 시작일 가능성이 높음
                if item['type'] == 'logical_row' and len(item['elements']) > 1:
                    table_data = []
                    col_count = len(item['elements'])
                    # 같은 열 개수를 가진 연속된 행들을 수집
                    while idx < len(rows) and rows[idx]['type'] == 'logical_row' and len(rows[idx]['elements']) == col_count:
                        table_data.append([clean_text(e['content']) for e in rows[idx]['elements']])
                        idx += 1
                    
                    if table_data:
                        # 데이터가 1줄 뿐이고 짧으면 표로 만들지 않음 (단순 나열)
                        if len(table_data) == 1 and sum(len(c) for c in table_data[0]) < 50:
                            processed_elements.append(" | ".join(table_data[0]) + "\n")
                        else:
                            # 마크다운 표 생성 (첫 행 헤더, 나머지 데이터)
                            if len(table_data) == 1:
                                headers = [f"Item {i+1}" for i in range(col_count)]
                                body = table_data
                            else:
                                headers = table_data[0]
                                body = table_data[1:]
                            
                            # 수동 표 생성 (표준 호환성 극대화)
                            hdr_str = "| " + " | ".join(headers) + " |"
                            sep_str = "| " + " | ".join([":---"] * len(headers)) + " |"
                            processed_elements.append("\n" + hdr_str + "\n" + sep_str)
                            for r_data in body:
                                processed_elements.append("| " + " | ".join(r_data) + " |")
                            processed_elements.append("\n")
                else:
                    if item['type'] == 'logical_row':
                        processed_elements.append(item['elements'][0]['content'] + "\n")
                    elif item['type'] == 'table':
                        # 진짜 표 객체 처리 (수동 구성)
                        # item['content']가 이미 tabulate된 문자열일 수 있으므로, 
                        # shape_type == TABLE 처리부에서 아예 데이터를 넘기도록 수정 필요 (아래에서 수정)
                        processed_elements.append("\n" + item['content'] + "\n")
                    else:
                        processed_elements.append("\n" + item['content'] + "\n")
                    idx += 1

            # 최종 출력
            for element in processed_elements:
                f.write(element + "\n")

    logger.info(f"변환 완료! 결과 파일: {output_md_path}")

# --- 실행 설정 ---
if __name__ == "__main__":
    input_pptx = os.path.join(os.path.dirname(__file__), "v3.pptx")
    output_md = os.path.join(os.path.dirname(__file__), "output.md")
    
    if os.path.exists(input_pptx):
        convert_pptx_to_md(input_pptx, output_md)
    else:
        logger.error(f"'{input_pptx}' 파일을 찾을 수 없습니다.")